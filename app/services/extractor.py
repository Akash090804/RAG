import requests
from bs4 import BeautifulSoup, Comment
import PyPDF2
from pptx import Presentation
import pandas as pd
from docx import Document
from io import BytesIO
from typing import Dict
import logging
from playwright.sync_api import sync_playwright
from urllib.parse import urlparse
import base64
import json
import re

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)



def fetch_content(url: str, headers: Dict[str, str] = None) -> bytes:
    """
    Fetch raw content from a URL. Uses Playwright for JavaScript-rendered HTML
    and the requests library for all other content types.
    """
    default_headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
    }
    if headers:
        default_headers.update(headers)
    try:
        head_response = requests.head(url, headers=default_headers, allow_redirects=True, timeout=10)
        content_type = head_response.headers.get('Content-Type', '').lower()
        if 'text/html' in content_type or not content_type:
            logger.info(f"Fetching HTML with Playwright: {url}")
            with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page()
                page.set_extra_http_headers(default_headers)
                page.goto(url, timeout=30000, wait_until='domcontentloaded')
                content = page.content().encode('utf-8')
                browser.close()
            return content
        else:
            logger.info(f"Fetching non-HTML content with requests: {url}")
            response = requests.get(url, headers=default_headers, timeout=30)
            response.raise_for_status()
            return response.content
    except requests.RequestException as e:
        logger.error(f"Requests failed to fetch URL {url}: {e}")
        raise ValueError(f"Failed to fetch URL with requests: {e}")
    except Exception as e:
        logger.error(f"Playwright failed for URL {url}: {e}")
        raise ValueError(f"Failed to fetch URL with Playwright: {e}")

def extract_text_from_content(content: bytes, url: str) -> str:
    """
    Extracts text from raw byte content based on the file type inferred from the URL.
    """
    parsed_url = urlparse(url)
    file_extension = parsed_url.path.split('.')[-1].lower() if '.' in parsed_url.path else ''
    logger.info(f"Extracting text for file type: {file_extension or 'html'}")
    try:
        if file_extension == 'pdf':
            with BytesIO(content) as f:
                reader = PyPDF2.PdfReader(f)
                return '\n'.join(page.extract_text() or '' for page in reader.pages)
        elif file_extension in ['ppt', 'pptx']:
            with BytesIO(content) as f:
                prs = Presentation(f)
                return '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text'))
        elif file_extension in ['xls', 'xlsx']:
            with BytesIO(content) as f:
                df = pd.read_excel(f, engine='openpyxl')
                return df.to_string(index=False)
        elif file_extension == 'docx':
            with BytesIO(content) as f:
                doc = Document(f)
                return '\n'.join(para.text for para in doc.paragraphs)
        elif file_extension in ['txt', 'csv', 'json']:
            return content.decode('utf-8', errors='ignore')
        else:
            soup = BeautifulSoup(content, 'lxml')
            text = soup.get_text(separator='\n', strip=True)
            hidden_texts = [el.get_text(separator='\n', strip=True) for el in soup.find_all(style=lambda v: v and 'display:none' in v)]
            comments = [str(comment).strip() for comment in soup.find_all(string=lambda text: isinstance(text, Comment))]
            full_text = '\n'.join([text] + hidden_texts + comments)
            logger.info(f"Extracted HTML text length: {len(full_text)} chars")
            return full_text
    except Exception as e:
        logger.error(f"Failed to extract text from content for {url}: {e}")
        raise ValueError(f"Text extraction failed: {e}")


def decode_jwt_payload(token: str) -> dict:
    """
    Decodes the payload of a JSON Web Token (JWT) without verification.
    This function ONLY deals with the token string.
    """
    try:
        _, payload_b64, _ = token.split('.')
        payload_b64 += '=' * (-len(payload_b64) % 4)
        payload_json = base64.urlsafe_b64decode(payload_b64).decode('utf-8')
        return json.loads(payload_json)
    except Exception as e:
        # This will now catch the "not enough values to unpack" error correctly
        logger.error(f"Failed to decode JWT token string: {e}")
        return {}

def solve_hidden_code_challenge(url: str) -> dict:
    """
    Solves the HackRx "Hidden Code" challenge by extracting the Challenge ID from the URL's JWT
    and finding the Completion Code from the dynamically rendered page content.
    """
    challenge_id = "Not found"
    completion_code = "Not found"


    try:
        # The logic that uses the 'url' variable is correctly placed here.
        token = url.split('/')[-1]
        jwt_payload = decode_jwt_payload(token)
        if jwt_payload: # Check if decoding was successful
             challenge_id = jwt_payload.get('challengeId', 'Challenge ID key not found')
             logger.info(f"Successfully decoded Challenge ID: {challenge_id}")
        else:
             challenge_id = "Failed to decode token from URL"
             logger.warning(challenge_id)
    except IndexError:
        # This handles cases where the URL doesn't have slashes, etc.
        logger.warning("Could not split URL to find token. This is normal for non-challenge URLs.")
        challenge_id = "Not applicable for this URL"

    # --- Step 2: Use Playwright to find the dynamically loaded Completion Code ---
    logger.info("Launching Playwright to find the completion code...")
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        try:
            page.goto(url, timeout=45000, wait_until='networkidle')
            completion_code_element = page.query_selector('#completionCode')
            if completion_code_element:
                completion_code = completion_code_element.get_attribute('value')
                logger.info(f"Found Completion Code in hidden input: {completion_code}")
            else:
                content = page.content()
                match = re.search(r'HR-\w{3}-\w{3}-\w{3}-\w{3}', content)
                if match:
                    completion_code = match.group(0)
                    logger.info(f"Found Completion Code via regex fallback: {completion_code}")
                else:
                    logger.warning("Could not find completion code on the page.")
        except Exception as e:
            logger.error(f"Playwright failed to retrieve the completion code: {e}")
            completion_code = f"Playwright Error: {e}"
        finally:
            browser.close()

    return {
        "challenge_id": challenge_id,
        "completion_code": completion_code
    }